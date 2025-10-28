from flask import Flask, request, send_file, jsonify
from flask_cors import CORS
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.enum.text import PP_ALIGN
import io
import traceback
import os

app = Flask(__name__)
CORS(app)

# Configuration
MAX_FILE_SIZE = 50 * 1024 * 1024  # 50MB
MAX_IMAGE_SIZE = 10 * 1024 * 1024  # 10MB
ALLOWED_IMAGE_EXTENSIONS = {'png', 'jpg', 'jpeg', 'gif', 'bmp'}

def validate_file_size(file_storage, max_size, file_type="File"):
    """Validate file size before processing"""
    file_storage.seek(0, os.SEEK_END)
    size = file_storage.tell()
    file_storage.seek(0)
    
    if size > max_size:
        size_mb = size / (1024 * 1024)
        max_mb = max_size / (1024 * 1024)
        raise ValueError(f"{file_type} too large: {size_mb:.2f}MB (max: {max_mb}MB)")
    
    return size

def validate_image_format(filename):
    """Validate image file extension"""
    if '.' not in filename:
        return False
    ext = filename.rsplit('.', 1)[1].lower()
    return ext in ALLOWED_IMAGE_EXTENSIONS

def find_layout_by_name(slide_master, layout_name):
    """Find layout by exact name match"""
    for layout in slide_master.slide_layouts:
        if layout.name == layout_name:
            return layout
    
    # If exact match fails, try case-insensitive partial match
    layout_name_lower = layout_name.lower()
    for layout in slide_master.slide_layouts:
        if layout_name_lower in layout.name.lower():
            return layout
    
    # Last resort: return first layout with content placeholders
    for layout in slide_master.slide_layouts:
        has_body = False
        for shape in layout.placeholders:
            if shape.placeholder_format.type in [PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT]:
                has_body = True
                break
        if has_body:
            return layout
    
    # Ultimate fallback: return first layout
    if len(slide_master.slide_layouts) > 0:
        return slide_master.slide_layouts[0]
    
    return None

@app.route('/api/add-slide', methods=['POST'])
def add_slide():
    try:
        # Validate file upload
        if 'file' not in request.files:
            return jsonify({'error': 'No PowerPoint file uploaded'}), 400
        
        file = request.files['file']
        
        if file.filename == '':
            return jsonify({'error': 'No file selected'}), 400
        
        if not file.filename.endswith('.pptx'):
            return jsonify({'error': 'File must be a .pptx PowerPoint file'}), 400
        
        # Validate file size
        try:
            validate_file_size(file, MAX_FILE_SIZE, "PowerPoint file")
        except ValueError as e:
            return jsonify({'error': str(e)}), 400
            
        # Get form data
        layout_name = request.form.get('layout', '').strip()
        title = request.form.get('title', '').strip()
        text = request.form.get('text', '').strip()
        image = request.files.get('image')
        
        # Validate layout name provided
        if not layout_name:
            return jsonify({'error': 'Layout name is required'}), 400
        
        # Validate required content
        if not text:
            return jsonify({'error': 'Content text is required'}), 400
        
        # Validate image if provided
        if image:
            if not validate_image_format(image.filename):
                return jsonify({'error': f'Invalid image format. Allowed: {", ".join(ALLOWED_IMAGE_EXTENSIONS)}'}), 400
            try:
                validate_file_size(image, MAX_IMAGE_SIZE, "Image")
            except ValueError as e:
                return jsonify({'error': str(e)}), 400
        
        # Validate position
        try:
            position = int(request.form.get('position', 0))
            if position < 0:
                return jsonify({'error': 'Position cannot be negative'}), 400
        except ValueError:
            return jsonify({'error': 'Position must be a valid number'}), 400

        # Load presentation
        try:
            prs = Presentation(file)
        except Exception as e:
            return jsonify({'error': f'Invalid or corrupted PowerPoint file: {str(e)}'}), 400

        # Find the layout by exact name
        slide_master = prs.slide_masters[0]
        
        if len(slide_master.slide_layouts) == 0:
            return jsonify({'error': 'Presentation has no slide layouts'}), 400
        
        slide_layout = find_layout_by_name(slide_master, layout_name)
        
        if not slide_layout:
            return jsonify({'error': f'Could not find layout: {layout_name}'}), 400

        print(f"\n=== Adding Slide ===")
        print(f"Requested Layout: {layout_name}")
        print(f"Using Layout: {slide_layout.name}")
        print(f"Title: {title}")
        print(f"Text: {text[:50]}..." if len(text) > 50 else f"Text: {text}")

        # Validate and adjust position
        total_slides = len(prs.slides)
        if position < 0:
            position = 0
        elif position > total_slides:
            position = total_slides

        # Create new slide
        new_slide = prs.slides.add_slide(slide_layout)
        
        # Move slide to correct position if needed
        if position < total_slides:
            try:
                # Use private API carefully - wrap in try/except for safety
                xml_slides = prs.slides._sldIdLst
                slides = list(xml_slides)
                
                # Validate we have the slide to move
                if len(slides) > 0:
                    moved_slide = slides.pop()
                    # Ensure position is within valid range
                    insert_pos = min(position, len(slides))
                    slides.insert(insert_pos, moved_slide)
                    xml_slides.clear()
                    for slide in slides:
                        xml_slides.append(slide)
                    print(f"✓ Slide inserted at position {insert_pos}")
                else:
                    print(f"! No slides to reposition")
            except AttributeError:
                print(f"! Slide positioning not supported in this version of python-pptx")
            except Exception as e:
                print(f"! Could not reposition slide: {str(e)}")
                # Don't fail the entire operation if positioning fails

        # SET TITLE
        title_set = False
        if title:
            try:
                if new_slide.shapes.title:
                    new_slide.shapes.title.text = title
                    print(f"✓ Title: '{title}'")
                    title_set = True
            except Exception as e:
                print(f"! Could not set title: {str(e)}")

        # SET CONTENT - SIMPLE AND CLEAN
        if text:
            # Try to find content placeholder
            content_added = False
            placeholder_found = False
            
            for shape in new_slide.placeholders:
                try:
                    phf = shape.placeholder_format
                    placeholder_found = True
                    
                    # Look for body/content placeholder
                    if phf.type in [PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT]:
                        # Found content placeholder - use it
                        tf = shape.text_frame
                        
                        # Validate text frame is writable
                        if not tf:
                            print(f"! Placeholder has no text frame")
                            continue
                        
                        tf.clear()
                        
                        # Add text with proper formatting
                        p = tf.paragraphs[0]
                        p.text = text
                        p.alignment = PP_ALIGN.LEFT
                        p.level = 0
                        
                        # Adjust text frame margins for better positioning
                        tf.margin_left = Inches(0.1)
                        tf.margin_top = Inches(0.1)
                        tf.word_wrap = True
                        
                        print(f"✓ Content: '{text[:50]}...' (in placeholder)")
                        content_added = True
                        break
                except AttributeError as e:
                    print(f"! Placeholder attribute error: {str(e)}")
                    continue
                except Exception as e:
                    print(f"! Error setting placeholder content: {str(e)}")
                    continue
            
            if not placeholder_found:
                print(f"⚠️ No placeholders found in slide layout")
            
            # If no placeholder found or failed, create text box
            if not content_added:
                try:
                    left = Inches(0.5)
                    top = Inches(1.8)
                    width = Inches(9)
                    height = Inches(5)
                    
                    txBox = new_slide.shapes.add_textbox(left, top, width, height)
                    tf = txBox.text_frame
                    p = tf.paragraphs[0]
                    p.text = text
                    p.alignment = PP_ALIGN.LEFT
                    
                    tf.word_wrap = True
                    tf.margin_left = Inches(0)
                    tf.margin_top = Inches(0)
                    
                    print(f"✓ Content: '{text[:50]}...' (in textbox fallback)")
                except Exception as e:
                    print(f"! Error creating textbox: {str(e)}")
                    raise ValueError("Could not add content to slide")

        # Add image if provided
        if image:
            image_added = False
            
            # Try to find image placeholder first
            for shape in new_slide.placeholders:
                try:
                    if shape.placeholder_format.type == PP_PLACEHOLDER.PICTURE:
                        # Found picture placeholder
                        shape.insert_picture(io.BytesIO(image.read()))
                        print(f"✓ Image added to placeholder")
                        image_added = True
                        break
                except Exception as e:
                    print(f"! Could not use picture placeholder: {str(e)}")
                    continue
            
            # If no placeholder, add as regular image
            if not image_added:
                try:
                    image.seek(0)  # Reset file pointer
                    image_stream = io.BytesIO(image.read())
                    
                    # Calculate position based on slide dimensions
                    left = Inches(6.5)
                    top = Inches(2)
                    max_height = Inches(4)
                    max_width = Inches(3)
                    
                    # Add image with size constraints
                    pic = new_slide.shapes.add_picture(image_stream, left, top, height=max_height)
                    
                    # Ensure width doesn't exceed bounds
                    if pic.width > max_width:
                        pic.width = max_width
                    
                    print(f"✓ Image added as shape")
                except Exception as e:
                    print(f"! Image error: {str(e)}")
                    # Don't fail the entire operation if image fails
                    pass

        print(f"=== Complete ===\n")
        
        # Save and return
        output = io.BytesIO()
        try:
            prs.save(output)
            output.seek(0)
        except Exception as e:
            print(f"ERROR saving presentation: {str(e)}")
            return jsonify({'error': f'Failed to save presentation: {str(e)}'}), 500

        return send_file(
            output, 
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
            as_attachment=True,
            download_name='presentation_modified.pptx'
        )
    
    except ValueError as e:
        # Validation errors
        return jsonify({'error': str(e)}), 400
        
    except Exception as e:
        print(f"ERROR: {str(e)}")
        print(traceback.format_exc())
        return jsonify({'error': f'Server error: {str(e)}'}), 500

@app.route('/api/health', methods=['GET'])
def health_check():
    return jsonify({'status': 'healthy', 'message': 'SlideMaker API is running'}), 200

@app.route('/api/get-layouts', methods=['POST'])
def get_layouts():
    """Get available slide layouts from uploaded presentation"""
    try:
        if 'file' not in request.files:
            return jsonify({'error': 'No PowerPoint file uploaded'}), 400
        
        file = request.files['file']
        
        if file.filename == '':
            return jsonify({'error': 'No file selected'}), 400
        
        if not file.filename.endswith('.pptx'):
            return jsonify({'error': 'File must be a .pptx PowerPoint file'}), 400
        
        # Load presentation
        try:
            prs = Presentation(file)
        except Exception as e:
            return jsonify({'error': f'Invalid PowerPoint file: {str(e)}'}), 400
        
        if len(prs.slide_masters) == 0:
            return jsonify({'error': 'No slide masters found in presentation'}), 400
        
        # Get layouts from first slide master
        slide_master = prs.slide_masters[0]
        layouts = []
        
        for idx, layout in enumerate(slide_master.slide_layouts):
            layouts.append({
                'index': idx,
                'name': layout.name,
                'type': 'custom'
            })
        
        return jsonify({
            'total_slides': len(prs.slides),
            'layouts': layouts,
            'message': f'Found {len(layouts)} layouts in presentation'
        }), 200
        
    except Exception as e:
        print(f"Error in get_layouts: {str(e)}")
        print(traceback.format_exc())
        return jsonify({'error': f'Server error: {str(e)}'}), 500

if __name__ == '__main__':
    app.run(debug=True, host='0.0.0.0', port=5000)
